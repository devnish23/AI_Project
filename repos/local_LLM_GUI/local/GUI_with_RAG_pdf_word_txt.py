import subprocess
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import messagebox, filedialog
import os
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd

def process_document(file_path):
    try:
        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        elif file_path.endswith('.pdf'):
            return extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            return extract_text_from_word(file_path)
        elif file_path.endswith('.pptx'):
            return extract_text_from_ppt(file_path)
        elif file_path.endswith('.xlsx'):
            return extract_text_from_excel(file_path)
        else:
            messagebox.showerror("Error", "Unsupported file format. Supported formats: .txt, .pdf, .docx, .pptx, .xlsx")
            return None
    except Exception as e:
        messagebox.showerror("Error", f"Failed to process document: {e}")
        return None

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_word(file_path):
    doc = Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

def upload_document():
    file_path = filedialog.askopenfilename(
        title="Select a Document",
        filetypes=(
            ("Supported Files", "*.txt *.pdf *.docx *.pptx *.xlsx"),
            ("Text Files", "*.txt"),
            ("PDF Files", "*.pdf"),
            ("Word Files", "*.docx"),
            ("PowerPoint Files", "*.pptx"),
            ("Excel Files", "*.xlsx"),
            ("All Files", "*.*")
        )
    )
    if file_path:
        document_content = process_document(file_path)
        if document_content:
            document_text.config(state="normal")
            document_text.delete("1.0", END)
            document_text.insert(END, document_content)
            document_text.config(state="disabled")
            messagebox.showinfo("Success", "Document uploaded successfully!")