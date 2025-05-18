import subprocess
import ttkbootstrap as ttk
from ttkbootstrap.constants import *
from tkinter import messagebox, filedialog
import os
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import ttkbootstrap as ttk  # Import ttkbootstrap for modern UI

def get_models():
    try:
        result = subprocess.run(['ollama', 'list'], capture_output=True, text=True, encoding='utf-8', check=True)
        models = result.stdout.strip().split('\n')[1:]  # Skip the header row
        return [line.split() for line in models]
    except subprocess.CalledProcessError as e:
        messagebox.showerror("Error", f"Failed to fetch models: {e}")
        return []

def refresh_models():
    models = get_models()
    for row in tree.get_children():
        tree.delete(row)
    for model in models:
        tree.insert("", "end", values=model)

def process_document(file_path):
    try:
        if file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        elif file_path.endswith('.pdf'):
            return extract_text_from_pdf(file_path)
        elif file_path.endswith('.docx'):
            return extract_text_from_word(file_path)
        elif file_path.endswith('.pptx'):
            return extract_text_from_ppt(file_path)
        elif file_path.endswith('.xlsx'):
            return extract_text_from_excel(file_path)
        else:
            messagebox.showerror("Error", "Unsupported file format. Supported formats: .txt, .pdf, .docx, .pptx, .xlsx")
            return None
    except Exception as e:
        messagebox.showerror("Error", f"Failed to process document: {e}")
        return None
def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_word(file_path):
    doc = Document(file_path)
    return "\n".join([paragraph.text for paragraph in doc.paragraphs])

def extract_text_from_ppt(file_path):
    presentation = Presentation(file_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string(index=False)

def clear_history():
    output_text.config(state="normal")  # Enable editing temporarily
    output_text.delete("1.0", END)  # Clear all text
    output_text.config(state="disabled")  # Disable editing to make it read-only

def upload_document():
    file_path = filedialog.askopenfilename(
        title="Select a Document",
        filetypes=(
            ("Supported Files", "*.txt *.pdf *.docx *.pptx *.xlsx"),
            ("Text Files", "*.txt"),
            ("PDF Files", "*.pdf"),
            ("Word Files", "*.docx"),
            ("PowerPoint Files", "*.pptx"),
            ("Excel Files", "*.xlsx"),
            ("All Files", "*.*")
        )
    )

    if file_path:
        document_content = process_document(file_path)
        if document_content:
            document_text.config(state="normal")
            document_text.delete("1.0", END)
            document_text.insert(END, document_content)
            document_text.config(state="disabled")
            messagebox.showinfo("Success", "Document uploaded successfully!")

def run_model_with_document():
    selected_item = tree.selection()
    if not selected_item:
        messagebox.showwarning("Warning", "Please select a model to run!")
        return

    model_name = tree.item(selected_item[0], 'values')[0]
    user_prompt = prompt_entry.get()

    if not user_prompt.strip():
        messagebox.showwarning("Warning", "Please enter a prompt to run the model!")
        return

    document_content = document_text.get("1.0", END).strip()
    if not document_content:
        messagebox.showwarning("Warning", "Please upload a document to use as context!")
        return

    combined_prompt = f"Document Context:\n{document_content}\n\nUser Prompt:\n{user_prompt}"

    try:
        # Run the command with the combined prompt
        result = subprocess.run(
            ['ollama', 'run', model_name],
            input=combined_prompt,  # Pass the combined prompt to the command
            capture_output=True,
            text=True,
            encoding='utf-8',
            check=True
        )
        # Append the command output to the output text area
        if result.stdout.strip():
            output_text.config(state="normal")  # Enable editing temporarily
            output_text.insert(END, f"Prompt: {user_prompt}\nDocument Context:\n{document_content}\nOutput:\n{result.stdout.strip()}\n{'-' * 50}\n")
            output_text.config(state="disabled")  # Disable editing to make it read-only
            output_text.see(END)  # Scroll to the end
        else:
            messagebox.showinfo("Output", "The model completed successfully but did not produce any output.")
    except subprocess.CalledProcessError as e:
        # Display the error output
        error_message = e.stderr.strip() if e.stderr else "An unknown error occurred."
        messagebox.showerror("Error", error_message)
    except Exception as ex:
        # Display unexpected exceptions
        messagebox.showerror("Error", f"An unexpected error occurred: {str(ex)}")

# Create a modern GUI using ttkbootstrap
app = ttk.Window(themename="superhero")  # Choose a modern theme
app.title("Ollama GUI")
app.geometry("800x600")
app.resizable(True, True)

# Frame for the main layout
frame = ttk.Frame(app, padding=10)
frame.pack(fill=BOTH, expand=True)

# Treeview to display models
tree_label = ttk.Label(frame, text="Available Models", font=("Helvetica", 12, "bold"))
tree_label.pack(anchor=W, pady=(0, 5))

tree = ttk.Treeview(frame, columns=("NAME", "ID", "SIZE", "MODIFIED"), show='headings', height=8, bootstyle=INFO)
tree.heading("NAME", text="NAME")
tree.heading("ID", text="ID")
tree.heading("SIZE", text="SIZE")
tree.heading("MODIFIED", text="MODIFIED")
tree.pack(fill=X, pady=(0, 10))

# Prompt input field
prompt_frame = ttk.Frame(frame)
prompt_frame.pack(fill=X, pady=(0, 10))

prompt_label = ttk.Label(prompt_frame, text="Enter Prompt:", font=("Helvetica", 10))
prompt_label.pack(side=LEFT, padx=(0, 5))

prompt_entry = ttk.Entry(prompt_frame, width=50)
prompt_entry.pack(side=LEFT, fill=X, expand=True)

# Buttons for actions
button_frame = ttk.Frame(frame)
button_frame.pack(fill=X, pady=(0, 10))

refresh_button = ttk.Button(button_frame, text="Refresh Models", bootstyle=SUCCESS, command=refresh_models)
refresh_button.pack(side=LEFT, padx=(0, 5))

run_button = ttk.Button(button_frame, text="Run Selected Model", bootstyle=PRIMARY, command=run_model_with_document)
run_button.pack(side=LEFT, padx=(0, 5))

upload_button = ttk.Button(button_frame, text="Upload Document", bootstyle=INFO, command=upload_document)
upload_button.pack(side=LEFT, padx=(0, 5))

clear_button = ttk.Button(button_frame, text="Clear History", bootstyle=WARNING, command=clear_history)
clear_button.pack(side=LEFT, padx=(0, 5))

# Document Text Area
document_frame = ttk.Labelframe(frame, text="Document Content", padding=10, bootstyle=LIGHT)
document_frame.pack(fill=BOTH, expand=True, pady=(0, 10))

document_text = ttk.Text(document_frame, wrap="word", height=10, state="disabled")
document_text.pack(side=LEFT, fill=BOTH, expand=True)

scrollbar_doc = ttk.Scrollbar(document_frame, orient="vertical", command=document_text.yview)
scrollbar_doc.pack(side=RIGHT, fill=Y)

document_text["yscrollcommand"] = scrollbar_doc.set

# Output Text Area with Scrollbar
output_frame = ttk.Labelframe(frame, text="Output", padding=10, bootstyle=DARK)
output_frame.pack(fill=BOTH, expand=True, pady=(0, 10))

output_text = ttk.Text(output_frame, wrap="word", height=15, state="disabled")
output_text.pack(side=LEFT, fill=BOTH, expand=True)

scrollbar = ttk.Scrollbar(output_frame, orient="vertical", command=output_text.yview)
scrollbar.pack(side=RIGHT, fill=Y)

output_text["yscrollcommand"] = scrollbar.set

# Initialize models
refresh_models()

# Run the main loop
app.mainloop()