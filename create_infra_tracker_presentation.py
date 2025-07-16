from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Slide content for each slide
slides = [
    {
        'title': 'Infra Tracker',
        'content': 'Infrastructure Management System\nComprehensive Solution for Application, Vendor, and Security Tracking'
    },
    {
        'title': 'Executive Summary',
        'content': '- Full-stack platform for infrastructure, patch, and vendor management\n- Real-time security advisory integration (Red Hat, etc.)\n- Automated reporting and scheduling\n- Modern, responsive UI for all stakeholders'
    },
    {
        'title': 'Key Features',
        'content': '- User authentication & role-based access\n- Application & vendor portal management\n- Bulk upload and file validation\n- Customizable reporting (Excel, PowerPoint)\n- Security risk tracking (CVEs, EOL dates)\n- Scheduling and notifications'
    },
    {
        'title': 'Architecture Overview',
        'content': '- Frontend: React, TypeScript, Tailwind CSS\n- Backend: Node.js, Express, MongoDB\n- Integration: Vendor APIs, Red Hat advisories\n- Deployment: Modular, scalable, cloud-ready'
    },
    {
        'title': 'System Diagram',
        'content': 'User → Frontend (React) → Backend (Express) → MongoDB\n                ↘ Vendor Portals (APIs)\n                ↘ File Storage\n\n(See documentation for detailed diagram)'
    },
    {
        'title': 'User Experience',
        'content': '- Clean, modern interface\n- Mobile-friendly, responsive design\n- Dark/light mode support\n- Interactive dashboards and charts'
    },
    {
        'title': 'Security & Compliance',
        'content': '- JWT authentication, password hashing\n- Role-based access control\n- Input validation, CORS, Helmet\n- Audit trails and logging'
    },
    {
        'title': 'Reporting & Analytics',
        'content': '- Generate Excel and PowerPoint reports\n- Customizable templates\n- Visual dashboards for management insights'
    },
    {
        'title': 'Automation & Scheduling',
        'content': '- Automated data fetching from vendor portals\n- Scheduled patch management and reporting\n- Email notifications (optional)'
    },
    {
        'title': 'Roadmap & Next Steps',
        'content': '- Expand vendor integrations (more platforms)\n- Advanced analytics and AI-driven insights\n- Enhanced notification and alerting\n- User feedback and continuous improvement'
    },
    {
        'title': 'Contact & Q&A',
        'content': '- Project team contacts\n- Questions & Discussion'
    },
]

def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for line in content.split('\n'):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.LEFT
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT

def add_architecture_diagram_slide(prs):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = 'Architecture Diagram'
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(2.0)
    height = Inches(0.8)
    # User
    user = slide.shapes.add_shape(1, left, top, width, height)
    user.text = 'User'
    # Frontend
    frontend = slide.shapes.add_shape(1, left + Inches(2.2), top, width, height)
    frontend.text = 'Frontend\n(React)'
    # Backend
    backend = slide.shapes.add_shape(1, left + Inches(4.4), top, width, height)
    backend.text = 'Backend\n(Express)'
    # MongoDB
    mongodb = slide.shapes.add_shape(1, left + Inches(6.6), top, width, height)
    mongodb.text = 'MongoDB'
    # Vendor Portals
    vendor = slide.shapes.add_shape(1, left + Inches(4.4), top + Inches(1.5), width, height)
    vendor.text = 'Vendor Portals'
    # File Storage
    storage = slide.shapes.add_shape(1, left + Inches(6.6), top + Inches(1.5), width, height)
    storage.text = 'File Storage'
    # Arrows
    slide.shapes.add_connector(2, user.left + user.width, user.top + user.height/2, frontend.left, frontend.top + frontend.height/2)
    slide.shapes.add_connector(2, frontend.left + frontend.width, frontend.top + frontend.height/2, backend.left, backend.top + backend.height/2)
    slide.shapes.add_connector(2, backend.left + backend.width, backend.top + backend.height/2, mongodb.left, mongodb.top + mongodb.height/2)
    slide.shapes.add_connector(2, backend.left + backend.width/2, backend.top + backend.height, vendor.left + vendor.width/2, vendor.top)
    slide.shapes.add_connector(2, backend.left + backend.width/2, backend.top + backend.height, storage.left + storage.width/2, storage.top)

prs = Presentation()

# Title slide (special layout)
title_slide_layout = prs.slide_layouts[0]
title_slide = prs.slides.add_slide(title_slide_layout)
title_slide.shapes.title.text = slides[0]['title']
# Robustly set subtitle text
for shape in title_slide.placeholders:
    if shape.placeholder_format.idx == 1:  # Subtitle placeholder
        shape.text_frame.text = slides[0]['content']
        break

# Architecture diagram slide
add_architecture_diagram_slide(prs)

# Other slides
for slide in slides[1:]:
    add_slide(prs, slide['title'], slide['content'])

prs.save('Infra_Tracker_Management_Presentation.pptx')
print('Presentation created: Infra_Tracker_Management_Presentation.pptx') 